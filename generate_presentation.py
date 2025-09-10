from pptx import Presentation
from pptx.util import Inches

# Create presentation
prs = Presentation()

# Slide 1: Introduction
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Introduction"
subtitle.text = "Name: Gautham T\nProject Title: English to Hindi Speech Translator\nTrack: Language Translation System"

# Slide 2: Problem Statement
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Problem Statement"
content.text = "The problem is the lack of easy-to-use tools for real-time speech translation from English to Hindi, especially for users who need quick and accurate translations with audio output.\n\nImportance: Enhances communication for non-English speakers, improves accessibility, and facilitates cross-language interactions in various contexts like education, business, and daily life."

# Slide 3: Project Outcome
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Project Outcome"
content.text = "The solution is a web application built with Flask that allows users to input English speech or text, translates it to Hindi using Google Translate API, and generates Hindi audio using text-to-speech synthesis.\n\nDemonstration: Users can record voice, type text, translate, and play audio output seamlessly."

# Slide 4: Objectives
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Objectives"
content.text = "Key goals of the project:\n- Provide accurate English to Hindi translation\n- Support both voice and text input\n- Generate high-quality Hindi audio output\n- Ensure responsive design for desktop and mobile\n- Maintain user-friendly interface"

# Slide 5: Methodology
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Methodology"
content.text = "Approach: Modular web application using Flask framework.\n\nDataset: Utilizes pre-trained models for speech recognition and synthesis.\n\nTools/Technologies:\n- Python, Flask, HTML/CSS/JavaScript\n- Google Translate API for translation\n- Vosk for speech-to-text\n- gTTS and pyttsx3 for text-to-speech\n- Web Speech API for browser-based voice input"

# Slide 6: Implementation
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Implementation"
content.text = "Developed the project by:\n- Setting up Flask application with routes for translation\n- Integrating STT module using Vosk model\n- Implementing translator module with Google Translate\n- Adding TTS module with IndicTrans and gTTS\n- Creating responsive HTML template with JavaScript for voice input\n- Handling audio file processing and base64 encoding for web transmission"

# Slide 7: Demo/Working
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Demo/Working"
content.text = "Demonstration of the project:\n1. Open the web app at http://localhost:5000\n2. Click 'Start Voice Input' and speak English\n3. Allow microphone permissions\n4. Click 'Translate' to get Hindi text\n5. Click 'Play Hindi Audio' to hear the translation\n\nAlternatively, type text and follow the same steps."

# Slide 8: Results & Observations
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Results & Observations"
content.text = "Achievements:\n- Successful integration of speech recognition, translation, and synthesis\n- Accurate translations with Google Translate\n- Clear audio output in Hindi\n\nPerformance:\n- Fast response times for text translation\n- Reliable voice input with browser API\n- Handles various accents and speech patterns\n- Responsive on multiple devices"

# Slide 9: Conclusion & Future Work
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion & Future Work"
content.text = "Summary of contributions:\n- Developed a functional speech translation web app\n- Demonstrated integration of multiple AI technologies\n- Provided a user-friendly solution for language barriers\n\nPossible improvements:\n- Add support for more languages\n- Implement offline capabilities\n- Enhance accuracy with custom models\n- Add conversation history and user accounts"

# Save the presentation
prs.save('project_presentation.pptx')